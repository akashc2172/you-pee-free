#!/usr/bin/env python3
"""Generate a scientific PPTX for the pediatric ureteral stent baseline/Tier A thesis deck."""

from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


SLATE = RGBColor(21, 36, 56)
TEAL = RGBColor(18, 114, 124)
GREEN = RGBColor(62, 128, 87)
AMBER = RGBColor(198, 138, 24)
RED = RGBColor(169, 64, 64)
LIGHT_BG = RGBColor(246, 248, 250)
MID_BG = RGBColor(233, 238, 243)
TEXT = RGBColor(23, 27, 31)
MUTED = RGBColor(92, 101, 110)
WHITE = RGBColor(255, 255, 255)


def set_background(slide) -> None:
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = WHITE


def add_header(slide, title: str, subtitle: str | None = None) -> None:
    set_background(slide)
    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.55)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = SLATE
    bar.line.color.rgb = SLATE

    title_box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(11.5), Inches(0.45))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = WHITE

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.75), Inches(12.0), Inches(0.35))
        tf = sub_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.color.rgb = MUTED


def add_footer(slide, index: int) -> None:
    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.45), Inches(7.05), Inches(12.4), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = MID_BG
    line.line.color.rgb = MID_BG

    foot = slide.shapes.add_textbox(Inches(0.55), Inches(7.08), Inches(12.0), Inches(0.22))
    tf = foot.text_frame
    tf.clear()
    left = tf.paragraphs[0]
    left.text = "Baseline COMSOL work only where completed; Tier A slides are clearly labeled as planned implementation."
    left.font.name = "Arial"
    left.font.size = Pt(8.5)
    left.font.color.rgb = MUTED

    num = slide.shapes.add_textbox(Inches(12.25), Inches(7.05), Inches(0.5), Inches(0.2))
    tf = num.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = str(index)
    p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Arial"
    p.font.size = Pt(9)
    p.font.color.rgb = MUTED


def add_badge(slide, x: float, y: float, w: float, text: str, color: RGBColor) -> None:
    box = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = color
    box.line.color.rgb = color
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = WHITE


def add_bullets(
    slide,
    items: list[str | tuple[str, int]],
    x: float,
    y: float,
    w: float,
    h: float,
    font_size: int = 18,
) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(0)
    tf.margin_bottom = Pt(0)
    first = True
    for item in items:
        text, level = item if isinstance(item, tuple) else (item, 0)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = text
        p.level = level
        p.font.name = "Arial"
        p.font.size = Pt(font_size if level == 0 else max(font_size - 2, 12))
        p.font.color.rgb = TEXT if level == 0 else MUTED
        p.space_after = Pt(5 if level == 0 else 2)
        p.bullet = True


def add_panel(slide, x: float, y: float, w: float, h: float, title: str, fill: RGBColor = LIGHT_BG):
    panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = fill
    panel.line.color.rgb = MID_BG
    if title:
        title_box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.08), Inches(w - 0.36), Inches(0.25))
        tf = title_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = title
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = SLATE
    return panel


def add_placeholder(slide, x: float, y: float, w: float, h: float, text: str) -> None:
    add_panel(slide, x, y, w, h, "Figure Placeholder", MID_BG)
    box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.5), Inches(w - 0.44), Inches(h - 0.65))
    tf = box.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.name = "Arial"
    p.font.size = Pt(14)
    p.font.italic = True
    p.font.color.rgb = MUTED


def add_text(slide, text: str, x: float, y: float, w: float, h: float, size: int = 16, color: RGBColor = TEXT, bold: bool = False, italic: bool = False, align: PP_ALIGN | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    if align is not None:
        p.alignment = align
    p.font.name = "Arial"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color


def draw_three_cylinder_schematic(slide, x: float, y: float, w: float, h: float) -> None:
    add_panel(slide, x, y, w, h, "3-Cylinder CFD Abstraction")
    region_y = y + 0.7
    region_h = 0.95
    segments = [
        (x + 0.25, 2.45, "Proximal cylinder\nrenal pelvis / kidney-side", RGBColor(203, 229, 255)),
        (x + 2.95, 4.2, "Middle cylinder\nureter shaft transport zone", RGBColor(210, 244, 229)),
        (x + 7.4, 2.45, "Distal cylinder\nbladder-side reservoir", RGBColor(255, 235, 206)),
    ]
    for seg_x, seg_w, label, color in segments:
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(seg_x),
            Inches(region_y),
            Inches(seg_w),
            Inches(region_h),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = WHITE
        tf = shape.text_frame
        tf.clear()
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = label
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = SLATE

    stent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x + 0.65), Inches(y + 1.05), Inches(8.55), Inches(0.18)
    )
    stent.fill.solid()
    stent.fill.fore_color.rgb = SLATE
    stent.line.color.rgb = SLATE

    add_text(
        slide,
        "Imported parametric stent passes through all three fluid regions",
        x + 0.8,
        y + 1.4,
        w - 1.6,
        0.35,
        size=12,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def add_table_like_slide(slide, headers: list[str], rows: list[list[str]], x: float, y: float, col_widths: list[float], row_h: float) -> None:
    current_x = x
    for header, width in zip(headers, col_widths):
        cell = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(current_x), Inches(y), Inches(width), Inches(row_h)
        )
        cell.fill.solid()
        cell.fill.fore_color.rgb = SLATE
        cell.line.color.rgb = WHITE
        tf = cell.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = header
        p.alignment = PP_ALIGN.CENTER
        p.font.name = "Arial"
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        current_x += width

    current_y = y + row_h
    for idx, row in enumerate(rows):
        current_x = x
        for value, width in zip(row, col_widths):
            cell = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(current_x), Inches(current_y), Inches(width), Inches(row_h)
            )
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG if idx % 2 == 0 else WHITE
            cell.line.color.rgb = MID_BG
            tf = cell.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = value
            p.font.name = "Arial"
            p.font.size = Pt(10)
            p.font.color.rgb = TEXT
            current_x += width
        current_y += row_h


def add_minimal_title(slide, title: str, subtitle: str | None = None) -> None:
    set_background(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1.0))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = SLATE
    p.alignment = PP_ALIGN.LEFT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.53), Inches(3.6), Inches(11.0), Inches(0.5))
        tf = sub_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = "Arial"
        p.font.size = Pt(18)
        p.font.color.rgb = MUTED
        p.alignment = PP_ALIGN.LEFT


def add_minimal_header(slide, title: str) -> None:
    set_background(slide)
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.0), Inches(0.6))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Arial"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = SLATE


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_minimal_title(
        slide, 
        "Pediatric Ureteral Stent Restoration", 
        "Baseline CFD Validation and Tier A Tissue-Context Planning"
    )
    add_footer(slide, 1)

    # 2. The Abstraction: 3-Cylinder Model
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_minimal_header(slide, "The 3-Cylinder CFD Abstraction")
    draw_three_cylinder_schematic(slide, 0.5, 1.2, 8.5, 2.5)
    add_bullets(
        slide,
        [
            "Simplifies complex anatomy into repeatable geometric regions.",
            "Kidney (Proximal) and Bladder (Distal) act as reservoir contexts.",
            "Middle Cylinder isolates the shaft transport zone for targeted study.",
        ],
        0.5, 4.0, 8.0, 3.0, font_size=20
    )
    add_placeholder(slide, 9.5, 1.2, 3.3, 5.0, "COMSOL Dumbbell\nGeometry Snapshot")
    add_footer(slide, 2)

    # 3. Current Baseline (Completed)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_minimal_header(slide, "Baseline Simulation Contract")
    add_badge(slide, 0.5, 1.1, 1.8, "COMPLETED", TEAL)
    add_bullets(
        slide,
        [
            ("Physics: Single-Phase Laminar Flow (spf)", 0),
            ("Boundary: ΔP = 490 Pa (Pressure-driven)", 0),
            ("Geometry: 11-parameter LHS sampled batch generation", 0),
            ("QC: Mass balance and mesh quality automated gating", 0),
        ],
        0.5, 1.8, 6.0, 4.0, font_size=22
    )
    add_placeholder(slide, 7.0, 1.2, 5.8, 5.0, "Baseline Velocity Field\n(Screenshot)")
    add_footer(slide, 3)

    # 4. Tier A: Peristaltic Planning
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_minimal_header(slide, "Tier A: Peristaltic Load Planning")
    add_badge(slide, 0.5, 1.1, 1.8, "PLANNED", AMBER)
    add_text(slide, "P_ext(z,t) = P0 + A sin(k z - ω t)", 0.5, 2.0, 6.0, 1.0, size=28, bold=True)
    add_bullets(
        slide,
        [
            "Target: Middle-cylinder radial loading only.",
            "Physiologic amplitude: 5–25 cm H20 scaling.",
            "Implementation: Prescribed geometry constriction phases.",
        ],
        0.5, 3.5, 6.0, 3.0, font_size=22
    )
    add_placeholder(slide, 7.0, 1.2, 5.8, 5.0, "Expected Tier A\nConstriction Schematic")
    add_footer(slide, 4)

    # 5. Scientific Guardrails
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_minimal_header(slide, "Evidence & Constraint Guardrails")
    headers = ["Domain", "Truth Claim", "Constraint"]
    rows = [
        ["Baseline", "Automation pipeline functional", "Rigid walls only"],
        ["Hierarchy", "Dumbbell is an abstraction", "Not anatomical surgery-planning"],
        ["Tier A", "Proposed loading model", "No fabricated deformable data"],
    ]
    add_table_like_slide(slide, headers, rows, 0.5, 1.5, [2.5, 5.0, 4.8], 0.8)
    add_text(slide, "The pipeline is designed for objective 100+ design campaigns, not patient-specific clinical diagnosis.", 0.5, 5.5, 12.0, 0.5, size=16, color=MUTED, italic=True)
    add_footer(slide, 5)

    # 6. Next Steps: Execution Roadmap
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_minimal_header(slide, "Next Steps & Roadmap")
    add_bullets(
        slide,
        [
            ("Harden mesh-clamping logic to stabilize batch 1 conversion.", 0),
            ("Execute 60-design LHS 'Stage 1' baseline campaign.", 0),
            ("Implement Tier A 'Single Phase' wave snapshot study.", 0),
            ("Retrain Bayesian Surrogate on loaded dataset.", 0),
        ],
        0.5, 1.5, 12.0, 5.0, font_size=24
    )
    add_footer(slide, 6)

    return prs


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate pediatric stent baseline/Tier A thesis presentation.")
    parser.add_argument(
        "--output",
        default="slides/pediatric_stent_comsol_thesis_tierA.pptx",
        help="Output PPTX path",
    )
    args = parser.parse_args()

    output = Path(args.output)
    output.parent.mkdir(parents=True, exist_ok=True)

    prs = build_presentation()
    prs.save(output)
    print(output)


if __name__ == "__main__":
    main()
